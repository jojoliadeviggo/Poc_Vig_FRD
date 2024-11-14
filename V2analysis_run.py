from pptx import Presentation
from keybert import KeyBERT
from transformers import pipeline
import torch
from dotenv import load_dotenv
import os
from mistralai import Mistral

load_dotenv()  # charger les variables d'environnement depuis .env

class PPTAnalyzer:
    def __init__(self):
        print("Initializing KeyBERT model...")
        self.keyword_model = KeyBERT(model="all-MiniLM-L6-v2")
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        self.mistral_client = Mistral(api_key=os.getenv("MISTRAL_API_KEY"))
        print(f"Using device: {self.device}")

    def extract_text(self, ppt_path):
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(ppt_path)
            text_content = []
        
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
        
            return " ".join(text_content)
            
        except Exception as e:
            print(f"Error extracting text: {str(e)}")
            raise

    def extract_keywords(self, text):
        """Extract keywords from text using KeyBERT"""
        print("\nExtracting keywords...")
        try:
            keywords = self.keyword_model.extract_keywords(
                text,
                keyphrase_ngram_range=(1, 2),  # Extract single words and pairs
                top_n=10,
                use_maxsum=True,
                nr_candidates=20
            )
            
            print("\nKeywords found:")
            for keyword, score in keywords:
                print(f"- {keyword}: {score:.3f}")
            
            return [kw for kw, _ in keywords]
            
        except Exception as e:
            print(f"Error extracting keywords: {str(e)}")
            return []

    def generate_summary(self, text, keywords):
        print("\nGenerating summary...")
        try:
            prompt = f"""Voici un texte extrait d'une présentation PowerPoint.
            Les mots-clés importants sont : {', '.join(keywords)}
            
            Texte : {text[:2000]}
            
            Tu dois générer un résumé concis (2-7 phrases) qui :
            1. Capture les points essentiels du document
            2. Intègre naturellement les mots-clés identifiés
            3. Est fidèle au contenu d'origine"""

            messages = [
                {
                    "role": "user",
                    "content": f"""Voici un texte extrait d'une présentation PowerPoint.
                    Les mots-clés importants sont : {', '.join(keywords)}
                    
                    Texte : {text[:2000]}
                    
                    Tu dois générer un résumé concis (2-7 phrases) qui :
                    1. Capture les points essentiels du document
                    2. Intègre naturellement les mots-clés identifiés
                    3. Est fidèle au contenu d'origine"""
                }
            ]


            response = self.mistral_client.chat.complete(
                model="mistral-tiny",  # ou "mistral-large-latest"
                messages=messages,
            )
            
            summary = response.choices[0].message.content
            print(f"\nGenerated summary: {summary}")
            return summary
            
        except Exception as e:
            print(f"Error generating summary: {str(e)}")
            return ""

    def analyze(self, ppt_path):
        """Main analysis function"""
        try:
            text = self.extract_text(ppt_path)
            print("\nExtracted text:")
            print(text)     
            if not text.strip():
                print("No text found in presentation")
                return {
                    "text_length": 0,
                    "keywords": [],
                    "summary" : ""
                }
            
            keywords = self.extract_keywords(text)
            summary = self.generate_summary(text, keywords)

            return {
                "text_length": len(text.split()),
                "keywords": keywords,
                "summary" : summary
            }
            
        except Exception as e:
            print(f"Analysis failed: {str(e)}")
            raise

# Example usage
if __name__ == "__main__":
    analyzer = PPTAnalyzer()
    result = analyzer.analyze("LIAaudeladubuzz.pptx")
    print("\nAnalysis results:", result)
