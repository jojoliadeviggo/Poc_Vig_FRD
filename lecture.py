from transformers import pipeline
import pptx
import os
import torch
from keybert import KeyBERT  
import traceback  # Pour avoir les erreurs détaillées

class PPTAnalyzer:
    def __init__(self, hf_token):
        self.summarizer = pipeline(
            "summarization",
            model="google/pegasus-x-large",  # nouveau modèle
            token=hf_token,
            device=0 if torch.cuda.is_available() else -1  # utilise GPU si disponible
        )

        # Commentez ou simplifiez le summarizer qui n'est pas nécessaire
        self.summarizer = None  # Temporairement désactivé

        # Gardez uniquement KeyBERT pour les mots-clés
        self.keyword_model = KeyBERT(model="all-MiniLM-L6-v2")
    
        self.classifier = pipeline(  # ça reste inchangé
            "zero-shot-classification",
            model="facebook/bart-large-mnli",
            token=hf_token
        )

    def clean_text(self, text):
        """Nettoie et prépare le texte pour l'analyse"""
        # Supprime les caractères spéciaux et les sauts de ligne multiples
        cleaned = ' '.join(text.split())
        # S'assure que le texte a une longueur minimale
        if len(cleaned) < 50:
            cleaned = cleaned + " " + cleaned
        return cleaned
    
    # Divise le texte en morceaux de taille similaire en respectant les phrases
    def split_into_chunks(self, text, chunk_size=8000):
        """Divise le texte en morceaux de taille similaire en respectant les phrases"""
        # S'assurer que le texte est assez long
        if len(text) < chunk_size:
            return [text]
        
        words = text.split()
        chunks = []
        current_chunk = []
        current_length = 0
    
        for word in words:
            current_chunk.append(word)
            current_length += len(word) + 1
        
            # Vérifier si on peut terminer le chunk
            if current_length >= chunk_size:
                # Chercher un point de coupure naturel
                chunk_text = ' '.join(current_chunk)
                last_period = max([chunk_text.rfind('.'), chunk_text.rfind('!'), chunk_text.rfind('?')])
            
                if last_period != -1:
                    # Couper à la dernière phrase
                    final_chunk = chunk_text[:last_period + 1].strip()
                    remainder = chunk_text[last_period + 1:].strip()
                
                    if len(final_chunk) >= 50:  # Vérifier taille minimale
                        chunks.append(final_chunk)
                        current_chunk = remainder.split() if remainder else []
                        current_length = len(remainder)
                else:
                    # Si pas de point trouvé, couper au dernier mot complet
                    chunks.append(chunk_text)
                    current_chunk = []
                    current_length = 0
    
        # Ajouter le dernier chunk s'il reste du texte
        final_text = ' '.join(current_chunk).strip()
        if final_text and len(final_text) >= 50:
            chunks.append(final_text)
        elif final_text and chunks:
            # Si le dernier morceau est trop court, l'ajouter au précédent
            chunks[-1] = chunks[-1] + ' ' + final_text
        
        return [chunk for chunk in chunks if len(chunk) >= 50]

    def extract_text(self, ppt_path):
        try:
            prs = pptx.Presentation(ppt_path)
            all_text = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    all_text.append(" ".join(slide_text))
            
            # Joindre tout le texte
            full_text = " ".join(all_text)
            print(f"Texte extrait (premiers 100 caractères): {full_text[:100]}...")
            return full_text
            
        except Exception as e:
            print(f"Erreur dans extract_text: {str(e)}")
            raise Exception(f"Erreur lors de l'extraction du texte : {str(e)}")

    def generate_summary(self, text):
        return "Analyse du résumé temporairement désactivée"
        try:
            # Nettoyer et préparer le texte
            cleaned_text = self.clean_text(text)
            print(f"Longueur du texte nettoyé: {len(cleaned_text)} caractères")

            # Vérifier la longueur minimale
            if len(cleaned_text) < 50:
                return "Texte trop court pour générer un résumé."

            # Diviser en chunks si nécessaire
            chunks = self.split_into_chunks(cleaned_text)
            print(f"Nombre de chunks à analyser: {len(chunks)}")

            # Générer un résumé pour chaque chunk
            summaries = []
            for i, chunk in enumerate(chunks, 1):
                try:
                    print(f"Traitement du chunk {i}/{len(chunks)} (longueur: {len(chunk)})")
                    if len(chunk) < 50:
                        print(f"Chunk {i} ignoré car trop court")
                        continue
                    
                    chunk_summary = self.summarizer(
                        chunk,
                        max_length=300, # Augmenté pour plus de détails
                        min_length=100, # Augmenté pour éviter les résumés trop courts
                        do_sample=True, # Activer l'échantillonnage pour plus de variété
                        temperature=0.7,        # Contrôle la créativité du modèle
                        top_p=0.85,            # Noyau de probabilité pour le texte généré
                        truncation=True,
                        length_penalty=1.5,  # favorise davantage les résumés plus longs
                        repetition_penalty=1.2,  # Évite les répétitions
                        num_beams=4  # améliore la qualité mais ralentit un peu (5:Augmenté pour une meilleure exploration)
                    )
                    if chunk_summary and len(chunk_summary) > 0:
                        summaries.append(chunk_summary[0]['summary_text'])
                    else:
                        print(f"Pas de résumé généré pour le chunk {i}")
                    
                except Exception as chunk_error:
                    print(f"Erreur sur le chunk {i}: {str(chunk_error)}")
                    print(f"Contenu du chunk problématique: {chunk[:100]}...")
                    continue  # Passer au chunk suivant en cas d'erreur
        
            if not summaries:
                return "Impossible de générer un résumé."
            
            if len(summaries) == 1:
                return summaries[0]

            # Combiner les résumés si multiples
            try:
                final_text = " ".join(summaries)
                final_summary = self.summarizer(
                    final_text,
                    max_length=200,
                    min_length=50,
                    do_sample=True,
                    temperature=0.6,     # Plus conservateur pour le résumé final
                    top_p=0.9,
                    num_beams=5,
                    length_penalty=1.2,
                    repetition_penalty=1.2,
                    truncation=True
                )
        
                print("\n=== RÉSUMÉ FINAL ===")
                print(final_summary[0]['summary_text'])
                return self.clean_summary(final_summary[0]['summary_text'])
                return final_summary[0]['summary_text']
            except Exception as e:
                print("Erreur lors de la génération du résumé final, retour du premier résumé")
                return summaries[0]

        except Exception as e:
            print(f"Erreur détaillée dans generate_summary: {str(e)}")
            print(f"Texte problématique: {text[:200]}...")
            raise Exception(f"Erreur lors de la génération du résumé : {str(e)}")

    def extract_keywords(self, text):
        try:
            cleaned_text = self.clean_text(text)
            print(f"\nDébut extraction mots-clés - Longueur texte: {len(cleaned_text)}")
            print(f"Exemple de texte: {cleaned_text[:200]}...")  # voir le début du texte

            # Limiter la taille du texte pour l'analyse
            max_length = 5000  # On limite à 5000 caractères
            if len(cleaned_text) > max_length:
                cleaned_text = cleaned_text[:max_length]
                print(f"Texte tronqué à {max_length} caractères pour l'analyse")
        
            # Analyser chaque chunk
            all_keywords = []
            for chunk in chunks:
                # Extraire les mots-clés avec KeyBERT
                keywords = self.keyword_model.extract_keywords(
                    cleaned_text,
                    keyphrase_ngram_range=(1, 2),  # Extraire des mots simples et des paires de mots
                    stop_words=None, # Désactivons les stop words pour voir 'french',  # Utiliser les stop words français
                    top_n=30,  # Nombre de mots-clés par chunk
                    diversity=0.1,  # Assure une diversité dans les résultats
                    use_maxsum=True,
                    nr_candidates=50  # augmente le nombre de candidats considérés
                )

                print("\nRésultats bruts de KeyBERT:")
                for kw, score in keywords:
                    print(f"- {kw}: {score}")

                # Garder tous les mots-clés trouvés, sans filtre de score
                keywords_list = [kw for kw, score in keywords]
        
                print(f"\nMots-clés retenus: {keywords_list}")
        
                return keywords_list or ["Aucun mot-clé trouvé"]  # Retourne au moins un élément

                # Debug - voir ce que keybert trouve
                print(f"Mots-clés trouvés pour ce chunk:", keywords)

                # Ajouter seulement les mots-clés avec un score suffisant
                chunk_keywords = [keyword for keyword, score in keywords if score > 0.05]
                print(f"Mots-clés retenus pour ce chunk:", chunk_keywords)  # Debug
                all_keywords.extend(chunk_keywords)
        
            # Retourner les mots-clés uniques
            # Debug final
            print(f"Tous les mots-clés trouvés:", all_keywords)
            unique_keywords = list(set(all_keywords)) if all_keywords else []
            print(f"Mots-clés uniques finaux:", unique_keywords)
        
            return unique_keywords or []  # S'assure de toujours retourner une liste

        except Exception as e:
            print(f"\nErreur détaillée dans extract_keywords: {str(e)}")
            traceback.print_exc()  # Affiche la stack trace complète
            return ["Erreur extraction mots-clés"]

    def analyze_ppt(self, ppt_path):
        try:
            print(f"Analyse du fichier: {ppt_path}")
            if not os.path.exists(ppt_path):
                raise FileNotFoundError(f"Le fichier {ppt_path} n'existe pas")
            
            # Extraire le texte
            text = self.extract_text(ppt_path)
            if not text.strip():
                return {
                    "title": "Sans titre",
                    "summary": "Fonction de résumé désactivée",
                    "keywords": [],
                    "text_length": 0
                }
            
            # Récupérer le titre
            prs = pptx.Presentation(ppt_path)
            title = ""
            if prs.slides:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if hasattr(shape, "text"):
                        title = shape.text.strip()
                        if title:
                            break

            # Analyser le contenu
            #summary = self.generate_summary(text)
            keywords = self.extract_keywords(text)
            
            return {
                "title": title or "Sans titre",
                "summary": "Fonction de résumé désactivée",
                "keywords": keywords,
                "text_length": len(text.split())
            }
            
        except Exception as e:
            print(f"Erreur complète dans analyze_ppt: {str(e)}")
            raise Exception(f"Erreur lors de l'analyse du PPT : {str(e)}")
        
            # Ajouter cette nouvelle méthode dans la classe
    def clean_summary(self, summary):
        """Nettoie et améliore le résumé final"""
        # Supprimer les phrases redondantes
        sentences = set(summary.split('.'))
        cleaned = '. '.join(sentences)
        # Corriger la ponctuation
        cleaned = cleaned.replace('..', '.').replace('. .', '.')
        # S'assurer que le résumé se termine par un point
        if not cleaned.endswith('.'):
            cleaned += '.'
        return cleaned.strip()