from pptx import Presentation
from keybert import KeyBERT
from transformers import pipeline
import torch
from dotenv import load_dotenv
import os
from mistralai import Mistral
import time

load_dotenv()  # charger les variables d'environnement depuis .env

class PPTAnalyzer:
    def __init__(self):
        print("Initializing KeyBERT model...")
        self.keyword_model = KeyBERT(model="all-MiniLM-L6-v2")
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        self.mistral_client = Mistral(api_key=os.getenv("MISTRAL_API_KEY"))
        print(f"Using device: {self.device}")

    def extract_text(self, ppt_path):
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(ppt_path)
            text_content = []
        
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
        
            return " ".join(text_content)
            
        except Exception as e:
            print(f"Error extracting text: {str(e)}")
            raise

    def extract_title(self, ppt_path):
        """Extract title from the first slide of PowerPoint file"""
        try:
            prs = Presentation(ppt_path)
            if prs.slides:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if hasattr(shape, "text"):
                        title = shape.text.strip()
                        if title:
                            return title
            return "Sans titre"
        except Exception as e:
            print(f"Error extracting title: {str(e)}")
            return "Sans titre"

    def extract_keywords(self, text):
        """Extract keywords from text using KeyBERT"""
        print("\nExtracting keywords...")
        try:
            keywords = self.keyword_model.extract_keywords(
                text,
                keyphrase_ngram_range=(1, 2),  # Extract single words and pairs
                top_n=10,
                use_maxsum=True,
                nr_candidates=20
            )
            
            print("\nKeywords found:")
            for keyword, score in keywords:
                print(f"- {keyword}: {score:.3f}")
            
            return [kw for kw, _ in keywords]
            
        except Exception as e:
            print(f"Error extracting keywords: {str(e)}")
            return []

    def generate_summary(self, text, keywords):
        print("\nGenerating summary...")
        try:
            time.sleep(1)  # Pause d'1 seconde avant l'appel API
            prompt = f"""Voici un texte extrait d'une présentation PowerPoint.
            Les mots-clés importants sont : {', '.join(keywords)}
            
            Texte : {text[:2000]}
            
            Tu dois générer un résumé concis (2-7 phrases) qui :
            1. Capture les points essentiels du document
            2. Intègre naturellement les mots-clés identifiés
            3. Est fidèle au contenu d'origine"""

            messages = [
                {
                    "role": "user",
                    "content": f"""Voici un texte extrait d'une présentation PowerPoint.
                    Les mots-clés importants sont : {', '.join(keywords)}
                    
                    Texte : {text[:2000]}
                    
                    Tu dois générer un résumé concis (2-7 phrases) qui :
                    1. Capture les points essentiels du document
                    2. Intègre naturellement les mots-clés identifiés
                    3. Est fidèle au contenu d'origine"""
                }
            ]


            response = self.mistral_client.chat.complete(
                model="mistral-tiny",  # ou "mistral-large-latest"
                messages=messages,
            )
            
            summary = response.choices[0].message.content
            print(f"\nGenerated summary: {summary}")
            return summary
            
        except Exception as e:
            print(f"Error generating summary: {str(e)}")
            return ""
    
    def extract_table_of_contents(self, text):
        print("\nExtracting/generating table of contents...")
        try:
            time.sleep(1)  # Pause d'1 seconde avant l'appel API
            messages = [
                {
                    "role": "user",
                    "content": f"""Analyse ce texte extrait d'une présentation PowerPoint et:
                    1. Si tu trouves des slides titrées 'Sommaire', 'Plan' ou 'Table des matières', extrais leur contenu.
                    2. Sinon, si tu identifies une structure claire avec des sections marquées par des titres de slides, génère un sommaire basé sur cette structure.
                    3. Si aucune structure claire n'est identifiable, retourne une chaîne vide.

                    Texte : {text[:2000]}

                    Important : Ne génère pas de sommaire artificiel sans structure claire."""
                }
            ]

            response = self.mistral_client.chat.complete(
                model="mistral-tiny",
                messages=messages,
            )
            
            toc = response.choices[0].message.content
            return toc if toc and not toc.lower().startswith("le texte ne") else ""
            
        except Exception as e:
            print(f"Error extracting table of contents: {str(e)}")
            return ""

    def analyze(self, ppt_path):
        """Main analysis function"""
        try:
            # Obtenir et afficher le chemin absolu
            abs_path = os.path.abspath(ppt_path)
            print(f"\nLe document analysé est situé ici : {abs_path}")

            text = self.extract_text(ppt_path)
            title = self.extract_title(ppt_path) 

            if not text.strip():
                print("No text found in presentation")
                return {
                    "title": "Sans titre",
                    "text_length": 0,
                    "keywords": [],
                    "summary" : "",
                    "table_of_contents": "",
                }
            
            keywords = self.extract_keywords(text)
            summary = self.generate_summary(text, keywords)

            return {
                "title": title,
                "text_length": len(text.split()),
                "keywords": keywords,
                "summary" : summary,
                "table_of_contents": self.extract_table_of_contents(text)  
            }
            
        except Exception as e:
            print(f"Analysis failed: {str(e)}")
            raise

# Example usage
if __name__ == "__main__":
    analyzer = PPTAnalyzer()
    file = r"C:\Users\Joséphine Balland\SCriptsfaitmaison\POC_FRD_2\Charte IA_VF_-Copie.pptx"
    result = analyzer.analyze(file)
    print("\nAnalysis results:", result)
